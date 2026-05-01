"""
PPT生成服务模块
负责生成HTML和PPTX格式的PPT
"""

import io
from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from .llm_service import llm_service, web_search_service
from .logger import logger

JAVASCRIPT_CODE = """
let slideCurrent = 0;
let slideTotal = 7;

function showSlide(n) {
    const slides = document.querySelectorAll('.slide');
    if (!slides || slides.length === 0) return;
    slideCurrent = n;
    if (slideCurrent >= slideTotal) slideCurrent = 0;
    if (slideCurrent < 0) slideCurrent = slideTotal - 1;
    slides.forEach(function(slide) {
        if (slide) {
            slide.style.display = 'none';
            slide.classList.remove('active');
        }
    });
    if (slides[slideCurrent]) {
        slides[slideCurrent].style.display = 'flex';
        slides[slideCurrent].classList.add('active');
    }
    var counter = document.getElementById('slideCounter');
    if (counter) counter.textContent = (slideCurrent + 1) + ' / ' + slideTotal;
}

function nextSlide() { showSlide(slideCurrent + 1); }
function prevSlide() { showSlide(slideCurrent - 1); }

function initSlides(total) {
    slideTotal = total || 7;
    slideCurrent = 0;
    showSlide(0);
    document.addEventListener('keydown', function(e) {
        if (e.key === 'ArrowRight' || e.key === ' ') nextSlide();
        if (e.key === 'ArrowLeft') prevSlide();
    });
}

document.addEventListener('DOMContentLoaded', function() {
    const prevBtn = document.getElementById('prevBtn');
    const nextBtn = document.getElementById('nextBtn');
    if (prevBtn) prevBtn.addEventListener('click', prevSlide);
    if (nextBtn) nextBtn.addEventListener('click', nextSlide);
});
"""

class PPTService:
    """PPT生成服务"""

    def __init__(self):
        self.scenario_configs = {
            "general": {
                "color_scheme": "#2E86AB",
                "font_family": "Microsoft YaHei, Arial, sans-serif",
                "style_class": "general-theme"
            },
            "technology": {
                "color_scheme": "#3498DB",
                "font_family": "Consolas, monospace",
                "style_class": "tech-theme"
            },
            "business": {
                "color_scheme": "#27AE60",
                "font_family": "Microsoft YaHei, Arial, sans-serif",
                "style_class": "business-theme"
            },
            "education": {
                "color_scheme": "#9B59B6",
                "font_family": "Microsoft YaHei, Arial, sans-serif",
                "style_class": "education-theme"
            },
            "tourism": {
                "color_scheme": "#E67E22",
                "font_family": "Microsoft YaHei, Arial, sans-serif",
                "style_class": "tourism-theme"
            },
            "science": {
                "color_scheme": "#1ABC9C",
                "font_family": "Microsoft YaHei, Arial, sans-serif",
                "style_class": "science-theme"
            }
        }

    async def generate_html(self, outline: Dict[str, Any], scenario: str = "general", use_llm: bool = False, use_web_search: bool = False, task_id: str = "") -> str:
        """生成HTML格式的PPT

        Args:
            outline: 大纲数据
            scenario: 应用场景
            use_llm: 是否使用LLM生成内容
            use_web_search: 是否使用联网搜索增强
            task_id: 进度跟踪任务ID

        Returns:
            HTML字符串
        """
        from .progress import progress_tracker

        logger.info(f"开始生成HTML PPT: scenario={scenario}, use_llm={use_llm}, use_web_search={use_web_search}")

        config = self.scenario_configs.get(scenario, self.scenario_configs["general"])
        topic = outline.get("title", "")
        language = outline.get("metadata", {}).get("language", "zh")

        slides_html = []
        slide_id = 1
        total_slides_count = len(outline.get("slides", []))

        for slide_idx, slide in enumerate(outline.get("slides", [])):
            slide_type = slide.get("type", "content")
            title = slide.get("title", "")
            enriched_slide = slide.copy()

            if task_id and progress_tracker.is_cancelled(task_id):
                break

            if slide_type == "content" and use_llm and llm_service.is_configured:
                progress_msg = f"📝 生成第 {slide_idx+1}/{total_slides_count} 页: {title}"
                progress_pct = 30 + int(60 * slide_idx / max(total_slides_count, 1))
                progress_tracker.update(task_id, progress_pct, progress_msg, f"gen_slide_{slide_idx}")
                logger.info(f"使用LLM生成幻灯片内容: {title}")
                try:
                    import asyncio
                    content = await asyncio.to_thread(
                        llm_service.generate_content, title, slide_type, topic, language, use_web_search, task_id
                    )
                    enriched_slide["content"] = content
                    logger.info(f"LLM内容生成成功: {title}")
                except Exception as e:
                    logger.error(f"LLM内容生成失败: {e}")

            if slide_type == "content":
                expanded_slides = self._expand_slide_with_template(enriched_slide, topic)
                for exp_slide in expanded_slides:
                    exp_slide["id"] = slide_id
                    slides_html.append(self._generate_slide_html(exp_slide, config))
                    slide_id += 1
            else:
                enriched_slide["id"] = slide_id
                slides_html.append(self._generate_slide_html(enriched_slide, config))
                slide_id += 1

        progress_tracker.update(task_id, 95, "🎨 渲染最终PPT页面...", "render_html")

        total_slides = len(slides_html)
        css_styles = self._generate_css_styles(config)

        complete_html = f"""<!DOCTYPE html>
<html lang="{language}">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{outline.get('title', 'PPT')}</title>
    <style>
        {css_styles}
    </style>
</head>
<body>
    <div class="presentation-container">
        <div class="slides-wrapper">
            {''.join(slides_html)}
        </div>
        <div class="navigation">
            <button id="prevBtn">‹ 上一页</button>
            <span id="slideCounter">1 / {total_slides}</span>
            <button id="nextBtn">下一页 ›</button>
        </div>
    </div>

    <script>
        {JAVASCRIPT_CODE}
        initSlides({total_slides});
    </script>
</body>
</html>
        """
        logger.info(f"HTML PPT生成完成: {total_slides} 张幻灯片, {len(complete_html)} 字符")
        return complete_html

    def _expand_slide_with_template(self, slide: Dict[str, Any], _topic: str) -> List[Dict[str, Any]]:
        """使用模板扩展内容幻灯片"""
        title = slide.get("title", "")
        content = slide.get("content", "")
        style = slide.get("style", {})

        points = self._parse_bullet_points(content)
        slides = []
        points_per_slide = 3

        for i in range(0, len(points), points_per_slide):
            chunk = points[i:i + points_per_slide]
            if chunk:
                slides.append({
                    "type": "content",
                    "title": title,
                    "subtitle": f"({i // points_per_slide + 1})" if len(points) > points_per_slide else "",
                    "content": '\n'.join(chunk),
                    "style": style
                })

        if not slides:
            slides.append({
                "type": "content",
                "title": title,
                "subtitle": "",
                "content": content,
                "style": style
            })

        return slides

    def _parse_bullet_points(self, content: str) -> List[str]:
        """解析bullet points"""
        if not content:
            return []

        lines = content.split('\n')
        points = []
        for line in lines:
            line = line.strip()
            if line:
                if not line.startswith('•') and not line.startswith('-') and not line.startswith('*'):
                    line = '• ' + line
                points.append(line)

        return points if points else [content]

    def _render_content_html(self, content: str, content_style: str = "bullet") -> str:
        """将内容文本渲染为HTML，支持不同内容样式"""
        import re
        content = content.replace('\n', '<br>')

        content = re.sub(r'\*\*(.+?)\*\*', r'<strong class="highlight-keyword">\1</strong>', content)
        content = content.replace('**', '')

        if content_style == "numbered":
            lines = content.split('<br>')
            numbered = []
            idx = 1
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                line = re.sub(r'^[•\-\*]\s*', '', line)
                numbered.append(f'<div class="numbered-item"><span class="num">{idx}</span><span class="num-text">{line}</span></div>')
                idx += 1
            return '<div class="numbered-list">' + ''.join(numbered) + '</div>'

        if content_style == "highlight":
            lines = content.split('<br>')
            items = []
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                line = re.sub(r'^[•\-\*]\s*', '', line)
                items.append(f'<div class="highlight-item">{line}</div>')
            return '<div class="highlight-list">' + ''.join(items) + '</div>'

        lines = content.split('<br>')
        items = []
        for line in lines:
            line = line.strip()
            if not line:
                continue
            if not line.startswith('•') and not line.startswith('-') and not line.startswith('*'):
                line = '• ' + line
            items.append(f'<div class="bullet-item">{line}</div>')
        return '<div class="bullet-list">' + ''.join(items) + '</div>'

    @staticmethod
    def _clean_markdown(text: str) -> str:
        import re
        text = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', text)
        text = text.replace('**', '')
        return text

    def _generate_slide_html(self, slide: Dict[str, Any], _config: Dict[str, Any]) -> str:
        slide_type = slide.get("type", "content")
        title = slide.get("title", "")
        subtitle = slide.get("subtitle", "")
        content = slide.get("content", "")
        style = slide.get("style", {})
        layout = style.get("layout", "left")
        accent_color = style.get("accent_color", "")
        icon = style.get("icon", "")
        content_style = style.get("content_style", "bullet")

        accent_style = f'style="--accent:{accent_color}"' if accent_color else ''

        if slide_type == "title":
            icon_html = f'<div class="slide-icon">{icon}</div>' if icon else ''
            return f'<div class="slide title-slide" {accent_style}>{icon_html}<h1>{title}</h1><h3>{subtitle}</h3></div>'

        if slide_type == "agenda":
            safe_content = self._clean_markdown(content)
            icon_html = f'<div class="slide-icon">{icon}</div>' if icon else ''
            return f'<div class="slide agenda-slide" {accent_style}>{icon_html}<h2>{title}</h2><div class="content">{safe_content}</div></div>'

        if slide_type == "thankyou":
            icon_html = f'<div class="slide-icon">{icon}</div>' if icon else ''
            return f'<div class="slide thankyou-slide" {accent_style}>{icon_html}<h2>{title}</h2><h3>{subtitle}</h3></div>'

        content_html = self._render_content_html(content, content_style)
        icon_html = f'<div class="slide-icon">{icon}</div>' if icon else ''

        if layout == "centered":
            return f'''<div class="slide content-slide layout-centered" {accent_style}>
                {icon_html}
                <h2>{title}</h2>
                {f'<h3>{subtitle}</h3>' if subtitle else ''}
                <div class="content-centered">{content_html}</div>
            </div>'''

        if layout == "quote":
            first_line = content.split('\n')[0].strip().lstrip('•-* ') if content else title
            first_line = self._clean_markdown(first_line)
            return f'''<div class="slide content-slide layout-quote" {accent_style}>
                {icon_html}
                <div class="quote-mark">"</div>
                <blockquote class="quote-text">{first_line}</blockquote>
                <p class="quote-source">— {title}</p>
            </div>'''

        if layout == "two-column":
            lines = [l.strip().lstrip('•-* ') for l in content.split('\n') if l.strip()]
            lines = [self._clean_markdown(l) for l in lines]
            mid = (len(lines) + 1) // 2
            left_items = lines[:mid]
            right_items = lines[mid:]
            left_html = '<br>'.join(f'• {l}' for l in left_items)
            right_html = '<br>'.join(f'• {l}' for l in right_items)
            return f'''<div class="slide content-slide layout-two-column" {accent_style}>
                {icon_html}
                <h2>{title}</h2>
                <div class="two-column-container">
                    <div class="column">{left_html}</div>
                    <div class="column">{right_html}</div>
                </div>
            </div>'''

        if layout == "cards":
            lines = [l.strip().lstrip('•-* ') for l in content.split('\n') if l.strip()]
            cards_html = ""
            for i, line in enumerate(lines[:4]):
                card_num = i + 1
                line = self._clean_markdown(line)
                cards_html += f'''<div class="card">
                    <div class="card-number">{card_num}</div>
                    <div class="card-text">{line}</div>
                </div>'''
            return f'''<div class="slide content-slide layout-cards" {accent_style}>
                {icon_html}
                <h2>{title}</h2>
                <div class="cards-container">{cards_html}</div>
            </div>'''

        if layout == "timeline":
            lines = [l.strip().lstrip('•-* ') for l in content.split('\n') if l.strip()]
            items_html = ""
            for i, line in enumerate(lines[:5]):
                line = self._clean_markdown(line)
                items_html += f'''<div class="timeline-item">
                    <div class="timeline-dot"></div>
                    <div class="timeline-content"><strong>{i+1}.</strong> {line}</div>
                </div>'''
            return f'''<div class="slide content-slide layout-timeline" {accent_style}>
                {icon_html}
                <h2>{title}</h2>
                <div class="timeline-container">{items_html}</div>
            </div>'''

        return f'''<div class="slide content-slide layout-left" {accent_style}>
            {icon_html}
            <h2>{title}</h2>
            {f'<h3>{subtitle}</h3>' if subtitle else ''}
            <div class="content">{content_html}</div>
        </div>'''

    def _generate_css_styles(self, config: Dict[str, Any]) -> str:
        """生成CSS样式"""
        return f"""
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        body {{
            font-family: {config['font_family']};
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
            display: flex;
            align-items: center;
            justify-content: center;
        }}
        .presentation-container {{
            width: 1280px;
            height: 720px;
            background: white;
            border-radius: 15px;
            box-shadow: 0 20px 40px rgba(0,0,0,0.1);
            overflow: hidden;
            margin: 0 auto;
        }}
        .slides-wrapper {{
            position: relative;
            width: 100%;
            height: 100%;
            overflow: hidden;
        }}
        .slide {{
            display: none;
            padding: 80px;
            width: 100%;
            height: 100%;
            background: white;
            flex-direction: column;
        }}
        .slide.active {{ display: flex; }}
        .slide-icon {{
            font-size: 2.5em;
            margin-bottom: 10px;
        }}
        .slide.title-slide {{
            background: linear-gradient(135deg, var(--accent, {config['color_scheme']}) 0%, {config['color_scheme']} 100%);
            color: white;
            justify-content: center;
            align-items: center;
        }}
        .slide.title-slide h1 {{
            font-size: 3em;
            margin-bottom: 20px;
            color: white;
        }}
        .slide h1 {{ font-size: 2.5em; margin-bottom: 30px; color: var(--accent, {config['color_scheme']}); }}
        .slide h2 {{ font-size: 2em; margin-bottom: 20px; color: var(--accent, {config['color_scheme']}); }}
        .slide h3 {{ font-size: 1.2em; color: #666; margin-bottom: 15px; }}
        .slide .content {{
            font-size: 1.1em;
            line-height: 1.8;
            color: #333;
        }}
        .slide.agenda-slide {{ background: #f8f9fa; }}
        .slide.thankyou-slide {{
            background: linear-gradient(135deg, var(--accent, {config['color_scheme']}) 0%, #764ba2 100%);
            color: white;
            justify-content: center;
            align-items: center;
        }}
        .slide.thankyou-slide h2 {{ color: white; font-size: 3em; }}
        .slide.thankyou-slide h3 {{ color: rgba(255,255,255,0.9); }}

        .layout-centered {{
            justify-content: center;
            align-items: center;
            text-align: center;
        }}
        .layout-centered .content-centered {{
            font-size: 1.3em;
            line-height: 2;
            max-width: 900px;
        }}

        .layout-quote {{
            justify-content: center;
            align-items: center;
            text-align: center;
        }}
        .quote-mark {{
            font-size: 6em;
            color: var(--accent, {config['color_scheme']});
            opacity: 0.3;
            line-height: 1;
            font-family: Georgia, serif;
        }}
        .quote-text {{
            font-size: 1.8em;
            font-style: italic;
            color: #333;
            max-width: 900px;
            line-height: 1.6;
            margin: 10px 0 20px;
        }}
        .quote-source {{
            font-size: 1.1em;
            color: #888;
        }}

        .layout-two-column .two-column-container {{
            display: flex;
            gap: 60px;
            flex: 1;
        }}
        .layout-two-column .column {{
            flex: 1;
            font-size: 1.1em;
            line-height: 2;
            color: #333;
            padding: 20px;
            background: #f8f9fa;
            border-radius: 10px;
            border-left: 4px solid var(--accent, {config['color_scheme']});
        }}

        .layout-cards .cards-container {{
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 20px;
            flex: 1;
        }}
        .layout-cards .card {{
            background: #f8f9fa;
            border-radius: 12px;
            padding: 25px;
            border-top: 4px solid var(--accent, {config['color_scheme']});
            box-shadow: 0 2px 8px rgba(0,0,0,0.06);
        }}
        .layout-cards .card-number {{
            font-size: 2em;
            font-weight: bold;
            color: var(--accent, {config['color_scheme']});
            margin-bottom: 8px;
        }}
        .layout-cards .card-text {{
            font-size: 1em;
            line-height: 1.6;
            color: #444;
        }}

        .layout-timeline .timeline-container {{
            position: relative;
            padding-left: 40px;
            flex: 1;
        }}
        .layout-timeline .timeline-container::before {{
            content: '';
            position: absolute;
            left: 15px;
            top: 0;
            bottom: 0;
            width: 3px;
            background: var(--accent, {config['color_scheme']});
            opacity: 0.3;
        }}
        .layout-timeline .timeline-item {{
            display: flex;
            align-items: flex-start;
            margin-bottom: 25px;
            position: relative;
        }}
        .layout-timeline .timeline-dot {{
            width: 14px;
            height: 14px;
            border-radius: 50%;
            background: var(--accent, {config['color_scheme']});
            position: absolute;
            left: -33px;
            top: 5px;
            box-shadow: 0 0 0 4px rgba(0,0,0,0.05);
        }}
        .layout-timeline .timeline-content {{
            font-size: 1.1em;
            line-height: 1.6;
            color: #333;
        }}

        .bullet-list .bullet-item {{
            font-size: 1.1em;
            line-height: 1.8;
            color: #333;
            padding: 4px 0;
        }}
        .numbered-list .numbered-item {{
            display: flex;
            align-items: flex-start;
            gap: 15px;
            margin-bottom: 12px;
        }}
        .numbered-list .num {{
            flex-shrink: 0;
            width: 32px;
            height: 32px;
            border-radius: 50%;
            background: var(--accent, {config['color_scheme']});
            color: white;
            display: flex;
            align-items: center;
            justify-content: center;
            font-weight: bold;
            font-size: 0.9em;
        }}
        .numbered-list .num-text {{
            font-size: 1.1em;
            line-height: 1.6;
            color: #333;
            padding-top: 4px;
        }}
        .highlight-list .highlight-item {{
            font-size: 1.15em;
            line-height: 1.8;
            color: #333;
            padding: 8px 15px;
            margin-bottom: 8px;
            background: #f8f9fa;
            border-left: 4px solid var(--accent, {config['color_scheme']});
            border-radius: 0 6px 6px 0;
        }}
        .highlight-keyword {{
            color: var(--accent, {config['color_scheme']});
            font-weight: bold;
        }}

        .navigation {{
            position: fixed;
            bottom: 20px;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            align-items: center;
            gap: 20px;
            background: white;
            padding: 10px 30px;
            border-radius: 50px;
            box-shadow: 0 5px 15px rgba(0,0,0,0.2);
        }}
        .navigation button {{
            padding: 10px 20px;
            border: none;
            background: {config['color_scheme']};
            color: white;
            border-radius: 25px;
            cursor: pointer;
            font-size: 1em;
            transition: all 0.3s;
        }}
        .navigation button:hover {{
            background: #764ba2;
            transform: scale(1.05);
        }}
        .navigation span {{
            font-size: 1.1em;
            color: #333;
            font-weight: bold;
        }}
        """

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

    def _add_shape(self, slide, left, top, width, height, fill_color, border_color=None):
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    def _add_rounded_rect(self, slide, left, top, width, height, fill_color, border_color=None):
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    def _add_circle(self, slide, left, top, size, fill_color):
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def _add_text(self, slide, left, top, width, height, text, font_size=20, bold=False, color=RGBColor(51,51,51), alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        return tf

    def _clean_md(self, text):
        import re
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = text.replace('**', '')
        return text

    def _lighten_color(self, color: RGBColor, factor=0.85) -> RGBColor:
        r = int(color[0] + (255 - color[0]) * factor)
        g = int(color[1] + (255 - color[1]) * factor)
        b = int(color[2] + (255 - color[2]) * factor)
        return RGBColor(min(r, 255), min(g, 255), min(b, 255))

    def generate_pptx(self, outline: Dict[str, Any], scenario: str = "general", use_llm: bool = False, use_web_search: bool = False) -> bytes:
        logger.info(f"开始生成PPTX: scenario={scenario}, use_llm={use_llm}, use_web_search={use_web_search}")

        config = self.scenario_configs.get(scenario, self.scenario_configs["general"])
        default_color = self._hex_to_rgb(config["color_scheme"])
        topic = outline.get("title", "")
        language = outline.get("metadata", {}).get("language", "zh")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        for slide_data in outline.get("slides", []):
            slide_type = slide_data.get("type", "content")
            title = slide_data.get("title", "")
            subtitle = slide_data.get("subtitle", "")
            content = slide_data.get("content", "")
            style = slide_data.get("style", {})
            accent_color_str = style.get("accent_color", "")
            layout = style.get("layout", "left")
            icon = style.get("icon", "")

            if accent_color_str:
                try:
                    slide_color = self._hex_to_rgb(accent_color_str)
                except Exception:
                    slide_color = default_color
            else:
                slide_color = default_color

            if slide_type == "content" and use_llm and llm_service.is_configured:
                try:
                    content = llm_service.generate_content(title, slide_type, topic, language, use_web_search=use_web_search)
                except Exception as e:
                    logger.error(f"LLM内容生成失败: {e}")

            content = self._clean_md(content)
            slide = prs.slides.add_slide(blank_layout)

            if slide_type == "title":
                self._build_title_slide(slide, title, subtitle, slide_color, icon)
            elif slide_type == "agenda":
                self._build_agenda_slide(slide, title, content, slide_color, icon)
            elif slide_type == "thankyou":
                self._build_thankyou_slide(slide, title, subtitle, slide_color, icon)
            else:
                self._build_content_slide_ex(slide, title, subtitle, content, slide_color, layout, icon)

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        pptx_bytes = buffer.read()

        logger.info(f"PPTX生成完成: {len(prs.slides)} 张幻灯片, {len(pptx_bytes)} 字节")
        return pptx_bytes

    def _build_title_slide(self, slide, title, subtitle, color, icon=""):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

        self._add_shape(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7), self._lighten_color(color, 0.7))

        if icon:
            self._add_text(slide, Inches(5.5), Inches(1.2), Inches(2.333), Inches(1), icon, font_size=48, alignment=PP_ALIGN.CENTER, color=RGBColor(255,255,255))

        title_top = Inches(2.2) if icon else Inches(2.5)
        self._add_text(slide, Inches(1), title_top, Inches(11.333), Inches(2), title, font_size=44, bold=True, alignment=PP_ALIGN.CENTER, color=RGBColor(255,255,255))

        if subtitle:
            self._add_text(slide, Inches(1), Inches(4.5), Inches(11.333), Inches(1), subtitle, font_size=24, alignment=PP_ALIGN.CENTER, color=RGBColor(230,230,230))

    def _build_agenda_slide(self, slide, title, content, color, icon=""):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(250, 250, 250)

        self._add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), color)

        if icon:
            self._add_text(slide, Inches(0.6), Inches(0.3), Inches(1), Inches(0.8), icon, font_size=28, color=color)

        self._add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1), title, font_size=36, bold=True, color=color)

        lines = [l.strip().lstrip('•-* ') for l in content.split('\n') if l.strip()]
        for i, line in enumerate(lines[:8]):
            y = Inches(1.8 + i * 0.65)
            self._add_circle(slide, Inches(1.2), y + Inches(0.05), Inches(0.25), color)
            self._add_text(slide, Inches(1.7), y, Inches(10), Inches(0.6), line, font_size=22, color=RGBColor(51,51,51))

    def _build_thankyou_slide(self, slide, title, subtitle, color, icon=""):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

        self._add_shape(slide, Inches(5.5), Inches(1.5), Inches(2.333), Inches(0.05), RGBColor(255,255,255))

        if icon:
            self._add_text(slide, Inches(5.5), Inches(2), Inches(2.333), Inches(1), icon, font_size=48, alignment=PP_ALIGN.CENTER, color=RGBColor(255,255,255))

        self._add_text(slide, Inches(1), Inches(3.2), Inches(11.333), Inches(1.5), title, font_size=48, bold=True, alignment=PP_ALIGN.CENTER, color=RGBColor(255,255,255))

        if subtitle:
            self._add_text(slide, Inches(1), Inches(4.8), Inches(11.333), Inches(1), subtitle, font_size=24, alignment=PP_ALIGN.CENTER, color=RGBColor(230,230,230))

    def _build_content_slide_ex(self, slide, title, subtitle, content, color, layout, icon=""):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        self._add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), color)
        self._add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), color)

        if icon:
            self._add_text(slide, Inches(0.6), Inches(0.3), Inches(1), Inches(0.8), icon, font_size=24, color=color)

        self._add_shape(slide, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.9), color)

        self._add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9), title, font_size=32, bold=True, color=color)

        if subtitle:
            self._add_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5), subtitle, font_size=16, color=RGBColor(128,128,128))

        lines = [l.strip().lstrip('•-* ') for l in content.split('\n') if l.strip()]

        if layout == "centered":
            self._layout_centered(slide, lines, color)
        elif layout == "quote":
            self._layout_quote(slide, lines, title, color)
        elif layout == "two-column":
            self._layout_two_column(slide, lines, color)
        elif layout == "cards":
            self._layout_cards(slide, lines, color)
        elif layout == "timeline":
            self._layout_timeline(slide, lines, color)
        else:
            self._layout_left(slide, lines, color)

    def _layout_left(self, slide, lines, color):
        for i, line in enumerate(lines[:8]):
            y = Inches(2.0 + i * 0.6)
            self._add_circle(slide, Inches(1.0), y + Inches(0.08), Inches(0.2), color)
            self._add_text(slide, Inches(1.4), y, Inches(10.5), Inches(0.55), line, font_size=20, color=RGBColor(51,51,51))

    def _layout_centered(self, slide, lines, color):
        for i, line in enumerate(lines[:6]):
            y = Inches(2.2 + i * 0.7)
            self._add_text(slide, Inches(2), y, Inches(9.333), Inches(0.6), line, font_size=22, alignment=PP_ALIGN.CENTER, color=RGBColor(51,51,51))

    def _layout_quote(self, slide, lines, title, color):
        if lines:
            quote_text = lines[0]
            self._add_text(slide, Inches(2), Inches(1.8), Inches(1), Inches(1.5), "\u201C", font_size=72, bold=True, color=self._lighten_color(color, 0.6))
            self._add_text(slide, Inches(2.5), Inches(2.5), Inches(8.5), Inches(2.5), quote_text, font_size=26, bold=True, color=RGBColor(51,51,51))
            self._add_text(slide, Inches(2.5), Inches(5.2), Inches(8.5), Inches(0.6), f"\u2014 {title}", font_size=18, color=RGBColor(128,128,128))

    def _layout_two_column(self, slide, lines, color):
        mid = (len(lines) + 1) // 2
        left_lines = lines[:mid]
        right_lines = lines[mid:]

        self._add_rounded_rect(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.8), RGBColor(248,249,250), color)
        self._add_rounded_rect(slide, Inches(6.9), Inches(2.0), Inches(5.8), Inches(4.8), RGBColor(248,249,250), color)

        for i, line in enumerate(left_lines[:6]):
            y = Inches(2.3 + i * 0.6)
            self._add_text(slide, Inches(1.0), y, Inches(5.0), Inches(0.55), f"\u2022 {line}", font_size=18, color=RGBColor(51,51,51))

        for i, line in enumerate(right_lines[:6]):
            y = Inches(2.3 + i * 0.6)
            self._add_text(slide, Inches(7.3), y, Inches(5.0), Inches(0.55), f"\u2022 {line}", font_size=18, color=RGBColor(51,51,51))

    def _layout_cards(self, slide, lines, color):
        card_lines = lines[:4]
        card_width = Inches(2.8)
        card_height = Inches(3.5)
        start_x = Inches(0.6)
        gap = Inches(0.3)

        for i, line in enumerate(card_lines):
            x = start_x + i * (card_width + gap)
            y = Inches(2.2)

            self._add_rounded_rect(slide, x, y, card_width, card_height, RGBColor(248,249,250), color)

            circle = self._add_circle(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.5), color)
            circle_tf = circle.text_frame
            circle_tf.paragraphs[0].text = str(i + 1)
            circle_tf.paragraphs[0].font.size = Pt(18)
            circle_tf.paragraphs[0].font.bold = True
            circle_tf.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
            circle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            self._add_text(slide, x + Inches(0.2), y + Inches(1.1), card_width - Inches(0.4), Inches(2.2), line, font_size=16, color=RGBColor(51,51,51))

    def _layout_timeline(self, slide, lines, color):
        line_y = Inches(4.0)
        self._add_shape(slide, Inches(1.5), line_y, Inches(10.333), Inches(0.06), self._lighten_color(color, 0.6))

        item_count = min(len(lines), 5)
        spacing = Inches(10.333) / max(item_count, 1)

        for i, line in enumerate(lines[:5]):
            x = Inches(1.5) + spacing * i
            dot_x = x + spacing / 2 - Inches(0.15)

            self._add_circle(slide, dot_x, line_y - Inches(0.12), Inches(0.3), color)

            self._add_text(slide, x, Inches(2.5), spacing, Inches(0.5), str(i + 1), font_size=28, bold=True, alignment=PP_ALIGN.CENTER, color=color)

            text_x = x - Inches(0.2)
            text_width = spacing + Inches(0.4)
            self._add_text(slide, text_x, Inches(4.5), text_width, Inches(2.0), line, font_size=15, alignment=PP_ALIGN.CENTER, color=RGBColor(51,51,51))

ppt_service = PPTService()
